"""Extract structured text from .pptx using python-pptx (one item per slide)."""

from __future__ import annotations

import io
import re
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    parts: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        line = "".join(run.text for run in paragraph.runs).strip()
        if line:
            parts.append(line)
    return "\n".join(parts)


def _table_text(shape) -> str:
    if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        return ""
    rows: list[str] = []
    table = shape.table
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _collect_slide_text(slide: Slide) -> tuple[str, str]:
    """Return (body_text, notes_text) for one slide."""
    body_parts: list[str] = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            t = _table_text(shape)
        else:
            t = _shape_text(shape)
        if t:
            body_parts.append(t)

    body = "\n".join(body_parts).strip()
    notes = ""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    return body, notes


def _guess_title(body: str) -> str:
    if not body:
        return "Untitled"
    first_line = body.split("\n", 1)[0].strip()
    if len(first_line) <= 120:
        return first_line or "Untitled"
    return first_line[:120] + "..."


def _detect_language(text: str) -> str:
    if not text:
        return "en"
    zh = len(re.findall(r"[\u4e00-\u9fff]", text))
    en = len(re.findall(r"[A-Za-z]", text))
    if zh > 0 and en > zh:
        return "mixed"
    if zh > 0:
        return "zh"
    return "en"


def extract_pptx_bytes(data: bytes) -> dict[str, Any]:
    prs = Presentation(io.BytesIO(data))
    pages: list[dict[str, Any]] = []
    warnings: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        body, notes = _collect_slide_text(slide)
        if not body and not notes:
            warnings.append(f"slide {idx}: empty content (image-only?)")
        pages.append(
            {
                "page_index": idx,
                "page_title": _guess_title(body),
                "original_text": body,
                "slide_notes": notes,
                "detected_language": _detect_language(body + "\n" + notes),
            }
        )

    return {
        "success": True,
        "total_pages": len(pages),
        "pages": pages,
        "warnings": warnings,
        "raw_char_count": sum(len(p["original_text"]) + len(p["slide_notes"]) for p in pages),
        "split_method": "python-pptx",
    }
